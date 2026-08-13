import sys
import os

# Ensure UTF-8 console output for Windows
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Tufte / Professional Business Controller Palette)
    C_DARK_BG = RGBColor(31, 41, 55)       # #1F2937 (Header / Dark theme)
    C_CARD_BG = RGBColor(249, 250, 251)    # #F9FAFB
    C_TEXT_MAIN = RGBColor(17, 24, 39)     # #111827
    C_TEXT_MUTED = RGBColor(107, 114, 128) # #6B7280
    C_WHITE = RGBColor(255, 255, 255)
    C_RED = RGBColor(220, 38, 38)          # #DC2626
    C_RED_BG = RGBColor(254, 226, 226)     # #FEE2E2
    C_AMBER = RGBColor(217, 119, 6)        # #D97706
    C_AMBER_BG = RGBColor(254, 243, 199)   # #FEF3C7
    C_GREEN = RGBColor(5, 150, 105)        # #059669
    C_GREEN_BG = RGBColor(209, 250, 229)   # #D1FAE5
    C_BLUE = RGBColor(37, 99, 235)         # #2563EB
    C_BORDER = RGBColor(229, 231, 235)     # #E5E7EB

    blank_slide_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, subtitle_text):
        # Header banner
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = C_DARK_BG
        header_shape.line.fill.background()

        tf = header_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(209, 213, 219)

    def add_kpi_card(slide, left, top, width, height, label, value, subtext, value_color=C_TEXT_MAIN, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(2)
        else:
            card.line.color.rgb = C_BORDER
            card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)

        p = tf.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_MUTED

        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = value_color

        p3 = tf.add_paragraph()
        p3.text = subtext
        p3.font.size = Pt(9)
        p3.font.color.rgb = C_TEXT_MUTED

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_DARK_BG
    bg1.line.fill.background()

    tf1 = bg1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(1.0)
    tf1.margin_top = Inches(2.2)

    p = tf1.paragraphs[0]
    p.text = "North Sea Oil Platform Drill Tower Construction"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    p_sub = tf1.add_paragraph()
    p_sub.text = "Integrated Earned Value Management (EVM), Schedule Control & Risk Briefing"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(156, 163, 175)

    p_meta = tf1.add_paragraph()
    p_meta.text = "\nStatus Date: August 31, 2026 (Month 8 / Week 36)   |   Prepared By: Frank Ellingsen, Lead Project Controller"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = RGBColor(209, 213, 219)

    # =========================================================================
    # SLIDE 2: EXECUTIVE STEERING COMMITTEE STATUS BRIEFING
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "1. Executive Steering Committee Status Briefing", "Overall Project Health: 🚨 CRITICAL COST OVERRUN & SCHEDULE SLIPPAGE (RED)")

    # Status Banner
    banner = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = C_RED_BG
    banner.line.color.rgb = C_RED
    banner.line.width = Pt(1.5)
    tf_b = banner.text_frame
    tf_b.margin_left = Inches(0.2)
    tf_b.margin_top = Inches(0.1)
    p = tf_b.paragraphs[0]
    p.text = "🚨 OVERALL HEALTH: CRITICAL (RED)   |   BAC Budget Burn Out: Month 9 (Sep 2026)   |   Required Financing: +$8.91M"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_RED

    # KPI Row
    add_kpi_card(slide2, Inches(0.5), Inches(2.0), Inches(2.3), Inches(1.1), "Baseline Budget (BAC)", "$26.50M", "Approved baseline target")
    add_kpi_card(slide2, Inches(3.0), Inches(2.0), Inches(2.3), Inches(1.1), "Earned Value (EV)", "$17.54M", "66.2% Work Complete", C_GREEN)
    add_kpi_card(slide2, Inches(5.5), Inches(2.0), Inches(2.3), Inches(1.1), "Actual Spend (AC)", "$23.44M", "Spent to date", C_RED)
    add_kpi_card(slide2, Inches(8.0), Inches(2.0), Inches(2.3), Inches(1.1), "Cost Index (CPI)", "0.75", "Severe overrun ($0.75/$1)", C_RED)
    add_kpi_card(slide2, Inches(10.5), Inches(2.0), Inches(2.333), Inches(1.1), "Time Index (SPIt)", "0.925", "-18.2 Days Schedule Slip", C_AMBER)

    # 2 Column Body: Left Outturns / Right Actions
    # Table of Scenarios
    rows, cols = 4, 3
    left_table = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(3.3), Inches(5.9), Inches(1.8)).table
    headers = ["Scenario", "EAC Outturn", "Completion Date"]
    for i, h in enumerate(headers):
        cell = left_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    data_scenarios = [
        ["Deterministic Base", "$35,413,604", "Jan 31, 2027 (+31d)"],
        ["Monte Carlo P50", "$34,060,783", "Feb 27, 2027 (+57d)"],
        ["Monte Carlo P90", "$35,815,202", "Mar 14, 2027 (+72d)"]
    ]
    for row_idx, row_data in enumerate(data_scenarios, start=1):
        for col_idx, text in enumerate(row_data):
            cell = left_table.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            if row_idx == 3:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # Right Side: Steering Actions Box
    actions_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(3.3), Inches(6.133), Inches(3.7))
    actions_box.fill.solid()
    actions_box.fill.fore_color.rgb = C_CARD_BG
    actions_box.line.color.rgb = C_BORDER
    tf_act = actions_box.text_frame
    tf_act.margin_left = Inches(0.2)
    tf_act.margin_top = Inches(0.2)
    
    p = tf_act.paragraphs[0]
    p.text = "REQUIRED STEERING COMMITTEE ACTIONS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_MAIN

    actions_list = [
        "1. Authorize Overrun Credit Line: Secure +$8.91M financing before September 15, 2026 (Month 9 budget exhaustion).",
        "2. Cap Yard Rework Billing: Transition Egersund yard (WBS 1.3.2) to a fixed-fee labor cap for remaining mast assembly.",
        "3. Compress Offshore Hook-up: Recover 15 calendar days to preserve January 31, 2027 Commercial COD.",
        "4. Enforce Change Freeze: Freeze all pending baseline scope change requests across WBS 1.3 deliverables."
    ]
    for act in actions_list:
        p = tf_act.add_paragraph()
        p.text = act
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_TEXT_MAIN
        p.space_after = Pt(8)

    # =========================================================================
    # SLIDE 3: RUNWAY & BURN STATISTICS (BUDGET DEPLETION)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "2. Monthly Cash Burn Speed & Budget Depletion Forecast", "Exact BAC Budget Exhaustion Point: Month 9 (September 2026)")

    add_kpi_card(slide3, Inches(0.5), Inches(1.4), Inches(2.9), Inches(1.2), "Avg Monthly Cash Burn", "$2.93M / Mo", "Baseline Plan: $2.56M/Mo (+14.5%)", C_RED)
    add_kpi_card(slide3, Inches(3.6), Inches(1.4), Inches(2.9), Inches(1.2), "Remaining Capital ($BAC-AC)", "$3.06M", "$23.44M spent of $26.50M BAC", C_AMBER)
    add_kpi_card(slide3, Inches(6.7), Inches(1.4), Inches(2.9), Inches(1.2), "Budget Burn Out Month", "Month 9 (Sep 2026)", "100% of BAC budget exhausted", C_RED, C_RED)
    add_kpi_card(slide3, Inches(9.8), Inches(1.4), Inches(3.033), Inches(1.2), "Required Overrun Financing", "+$8,913,604", "Must be approved before Sep 15", C_RED)

    # Table of Monthly Cash Burn Speed
    burn_table_shape = slide3.shapes.add_table(10, 5, Inches(0.5), Inches(2.8), Inches(12.333), Inches(4.3)).table
    b_headers = ["Month", "Planned Monthly ($PV)", "Earned Monthly ($EV)", "Actual Monthly ($AC)", "Budget Status & Runway Note"]
    for i, h in enumerate(b_headers):
        cell = burn_table_shape.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    b_data = [
        ["Month 1 (Jan)", "$0.30M", "$0.30M", "$0.32M", "Engineering ramp-up"],
        ["Month 2 (Feb)", "$0.90M", "$0.84M", "$0.95M", "Detail steel design completed"],
        ["Month 3 (Mar)", "$2.70M", "$2.26M", "$2.90M", "Procurement placement delays"],
        ["Month 4 (Apr)", "$3.80M", "$3.52M", "$4.05M", "Verdal yard fabrication start"],
        ["Month 5 (May)", "$3.00M", "$2.90M", "$3.22M", "Sub-structure assembly"],
        ["Month 6 (Jun)", "$4.00M", "$2.75M", "$4.45M", "Egersund mast dimensional fitting rework"],
        ["Month 7 (Jul)", "$3.00M", "$2.62M", "$3.85M", "Double-shift NDT welding hours"],
        ["Month 8 (Aug)", "$2.80M", "$2.35M", "$3.70M", "Status Date: $23.44M cumulative spend"],
        ["Month 9 (Sep)", "$2.50M (FC)", "$2.10M (FC)", "$3.20M (FC)", "💥 100% BAC Budget Exhausted (Sep 2026)"]
    ]
    for r_idx, r_data in enumerate(b_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = burn_table_shape.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            if r_idx == 9:
                p.font.bold = True
                p.font.color.rgb = C_RED
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_RED_BG

    # =========================================================================
    # SLIDE 4: COST VARIANCE WATERFALL BRIDGE
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "3. Cost Variance Waterfall Bridge ($BAC → EAC Outturn)", "Step-by-Step Cost Variance Drivers from Baseline to Outturn Forecast")

    wf_table = slide4.shapes.add_table(11, 4, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    wf_headers = ["WBS / Step Description", "Variance Amount", "Impact Category", "Driver & Root Cause Analysis"]
    for i, h in enumerate(wf_headers):
        cell = wf_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_WHITE

    wf_data = [
        ["Approved Baseline Budget ($BAC)", "$26.50M", "Baseline", "Original contract baseline target"],
        ["WBS 1.1 Engineering", "+$0.30M", "Scope Rework", "Minor drafting revision hours"],
        ["WBS 1.2 Procurement", "+$0.60M", "Material Inflation", "High-grade tubular steel price surge"],
        ["WBS 1.3.1 Verdal Yard Fabrication", "+$0.20M", "Yard Overtime", "Welder shift overtime"],
        ["WBS 1.3.2 Egersund Mast Assembly", "+$2.40M", "Primary Rework Driver", "Pipe fitting misalignment & 24/7 NDT welding"],
        ["WBS 1.4.1 Heavy Lift Vessel Mobilization", "+$1.30M", "Marine Standby", "Autumn sea-state vessel standby daily rates"],
        ["WBS 1.4.2 Topside Mating", "+$0.10M", "Offshore Fit", "Offshore trial fitting hours"],
        ["WBS 1.5 Commissioning & Hook-up", "+$0.20M", "System Test", "Pre-commissioning loops"],
        ["Time Delay Overhead Spread (SPIt)", "+$3.81M", "Schedule Overhead", "Extending execution past Dec 31 spreads PMO & site overhead"],
        ["Final Outturn Forecast ($EAC)", "$35.41M", "Total EAC", "Total Outturn Deficit: +$8.91M VAC (-33.6%)"]
    ]
    for r_idx, r_data in enumerate(wf_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = wf_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            if r_idx in [5, 9]:
                p.font.bold = True
                p.font.color.rgb = C_RED
            elif r_idx == 10:
                p.font.bold = True
                p.font.color.rgb = C_TEXT_MAIN
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BORDER

    # =========================================================================
    # SLIDE 5: S-CURVE PERFORMANCE & VARIANCE ANALYSIS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "4. Performance Measurement Baseline (S-Curves) & Variances", "Direct Variance Metrics & Time Slippage Analysis at Status Date (Aug 31, 2026)")

    add_kpi_card(slide5, Inches(0.5), Inches(1.4), Inches(2.9), Inches(1.2), "Cost Variance (CV = EV - AC)", "-$5,900,000", "-33.6% Cost Overrun", C_RED)
    add_kpi_card(slide5, Inches(3.6), Inches(1.4), Inches(2.9), Inches(1.2), "Schedule Variance (SV = EV - PV)", "-$2,960,000", "-14.4% Value Delay", C_AMBER)
    add_kpi_card(slide5, Inches(6.7), Inches(1.4), Inches(2.9), Inches(1.2), "Earned Schedule (ES)", "7.40 Months", "Actual Time Elapsed: 8.0M", C_AMBER)
    add_kpi_card(slide5, Inches(9.8), Inches(1.4), Inches(3.033), Inches(1.2), "Time Variance (SVt)", "-18.2 Days", "Behind Baseline Plan", C_RED)

    scurve_table = slide5.shapes.add_table(7, 5, Inches(0.5), Inches(2.8), Inches(12.333), Inches(4.3)).table
    sc_headers = ["Performance Metric", "Formula / Standard", "Current Value", "Threshold Status", "Operational Meaning"]
    for i, h in enumerate(sc_headers):
        cell = scurve_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    sc_data = [
        ["Planned Value (PV)", "Baseline Schedule S-Curve", "$20,500,000", "Baseline Target", "77.36% of work scope scheduled"],
        ["Earned Value (EV)", "Physical % Complete × BAC", "$17,540,000", "In Progress", "66.19% of work scope accomplished"],
        ["Actual Spend (AC)", "Cumulative General Ledger", "$23,440,000", "🚨 Exceeds EV", "Actual funds disbursed"],
        ["Cost Performance Index (CPI)", "EV / AC", "0.7483", "🚨 Red (< 0.90)", "$0.75 earned value generated for every $1.00 spent"],
        ["Schedule Index (SPI)", "EV / PV", "0.8556", "🟡 Amber (< 0.95)", "Scope earned velocity lags baseline schedule"],
        ["TCPI (Target BAC)", "(BAC - EV) / (BAC - AC)", "6.07", "🚨 Unviable (> 1.10)", "Impossible cost efficiency required to hit $26.50M BAC"]
    ]
    for r_idx, r_data in enumerate(sc_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = scurve_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            if r_idx in [3, 4, 6] and c_idx in [2, 3]:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # =========================================================================
    # SLIDE 6: EXECUTIVE VARIANCE EXPLANATIONS & CORRECTIVE ACTION AUDIT
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "5. Executive Variance Explanations & Corrective Actions", "Root Cause Analysis & Accountable CAM Response Actions")

    var_table = slide6.shapes.add_table(6, 5, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    v_headers = ["Variance Indicator", "Status Metric", "Root Cause Explanation", "Actionable Next Steps", "Owner"]
    for i, h in enumerate(v_headers):
        cell = var_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    v_data = [
        ["Cost Variance (CV)", "-$5,900,000\n(CPI = 0.75)", "Structural fitting rework & premium overtime at Verdal and Egersund yards.", "Enforce Design Change Freeze on WBS 1.3. Convert open T&M subcontracts to fixed unit-rate capped milestones.", "O. Eriksen"],
        ["Outturn Deficit (VAC)", "-$8,913,604\n(EAC: $35.41M)", "Linear extrapolation of poor cost efficiency (CPI=0.75) yields 33.6% overrun.", "Submit formal EAC re-baseline proposal. Request $3.50M drawdown from Management Reserve.", "Frank Ellingsen"],
        ["TCPI Target Efficiency", "6.07 (BAC)\n(Unviable)", "Original BAC ceiling requires impossible 607% cost efficiency for remaining work.", "Formally transition target to TCPI_EAC = 1.00 based on approved $35.41M EAC.", "Frank Ellingsen"],
        ["Schedule Delay (SV)", "-$2,960,000\n(+31 Days)", "Tubular steel procurement delays cascaded to Verdal and Egersund yards.", "Fast-track Egersund rigging with parallel work-fronts. Negotiate 10-day vessel window extension.", "M. Berg / K. Solberg"],
        ["WBS 1.3.2 Mast Assembly", "-$2,330,000\n(CPI = 0.57)", "Egersund rigging yard productivity bottleneck and excessive NDT rework.", "Deploy resident Project Controller to Egersund. Audit daily timecards & weld gates.", "O. Eriksen"]
    ]
    for r_idx, r_data in enumerate(v_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = var_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # =========================================================================
    # SLIDE 7: CFO FINANCIAL OUTTURN & CONTROL ACCOUNT FORECAST
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "6. CFO Outturn Forecast & Control Account Financials", "Control Account Level EAC Outturns, Cost Variances & Remaining ETC Liquidity")

    cfo_table = slide7.shapes.add_table(5, 8, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    c_headers = ["WBS Code", "Control Account", "BAC Budget", "Earned (EV)", "Actual (AC)", "Cost Var (CV)", "CPI", "Outturn EAC"]
    for i, h in enumerate(c_headers):
        cell = cfo_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    c_data = [
        ["1.1.1", "Structural Steel Detail Engineering", "$1,200,000", "$1,200,000", "$1,290,000", "-$90,000", "0.93", "$1,290,000"],
        ["1.2.1", "High-Grade Tubular Steel Procurement", "$3,500,000", "$3,500,000", "$3,770,000", "-$270,000", "0.93", "$3,770,000"],
        ["1.3.1", "Yard Sub-Structure Fabrication (Verdal)", "$4,000,000", "$3,600,000", "$4,450,000", "-$850,000", "0.81", "$4,944,000"],
        ["1.3.2", "Derrick Tower Mast Assembly (Egersund)", "$4,800,000", "$3,120,000", "$5,450,000", "-$2,330,000", "0.57", "$8,385,000"]
    ]
    for r_idx, r_data in enumerate(c_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = cfo_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            if c_idx in [5, 7] and r_idx == 4:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # =========================================================================
    # SLIDE 8: GANTT CRITICAL PATH & PREDECESSOR LOGIC
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "7. Offshore EPC Gantt Schedule & Critical Path Logic", "10 Major Tasks, Predecessor Links, Critical Path (8 Tasks) & Key Milestones")

    gantt_table = slide8.shapes.add_table(11, 7, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    g_headers = ["Task ID", "WBS", "Task Deliverable / Milestone Name", "Predecessor", "% Complete", "Forecast Window", "Critical Path"]
    for i, h in enumerate(g_headers):
        cell = gantt_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    g_data = [
        ["T101", "1.1.1", "Structural Steel Detail Engineering", "-", "100%", "Jan 05 - Mar 15", "Yes"],
        ["M1", "1.1.0", "◆ Milestone: Engineering AFC Gate Review", "T101 (FS)", "100%", "Achieved Mar 15", "Achieved"],
        ["T102", "1.1.2", "Piping & Drilling Package Design", "T101 (FS)", "100%", "Feb 01 - Apr 15", "No"],
        ["T103", "1.2.1", "High-Grade Tubular Steel Procurement", "T101 (+5D)", "100%", "Mar 15 - May 31", "Yes (+31D)"],
        ["T104", "1.2.2", "Mud Pumps & Top Drive Equipment", "T102 (FS)", "85%", "Apr 01 - Jul 31", "No"],
        ["T105", "1.3.1", "Yard Sub-Structure Fabrication (Verdal)", "T103 (FS)", "90%", "Jun 01 - Aug 31", "Yes"],
        ["M3", "1.3.0", "◆ Milestone: Verdal Sub-Structure Handover", "T105 (FS)", "90%", "Aug 31", "Active"],
        ["T106", "1.3.2", "Derrick Tower Mast Assembly (Egersund)", "T105 (FS)", "65%", "Jul 01 - Sep 30", "Yes (-30D Float)"],
        ["M4", "1.3.3", "◆ Milestone: Derrick Mast Ready for Heavy Lift", "T106 (FS)", "65%", "Sep 30 Target", "Critical Gate"],
        ["T107", "1.4.1", "Heavy Lift Vessel Mobilization (Heerema)", "T106 (FS)", "30%", "Sep 15 - Oct 15", "Yes"]
    ]
    for r_idx, r_data in enumerate(g_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = gantt_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            if c_idx == 6 and val.startswith("Yes"):
                p.font.bold = True
                p.font.color.rgb = C_RED

    # =========================================================================
    # SLIDE 9: EPC EXECUTIVE RISK MATRIX & HEATMAP
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "8. EPC Executive Risk Matrix & Mitigation Register", "5x5 Quantitative Probability vs. Impact Heatmap & Primary Risk Exposures")

    risk_table = slide9.shapes.add_table(6, 6, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    r_headers = ["Risk ID", "Risk Event Description", "WBS", "Owner", "Exposure ($)", "Mitigation Strategy"]
    for i, h in enumerate(r_headers):
        cell = risk_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    r_data = [
        ["R01 (Score 20)", "Egersund Mast Assembly Dimensional Out-of-Tolerance", "1.3.2", "O. Eriksen", "$2,400,000", "24/7 NDT welding specialists; fixed-fee labor cap."],
        ["R02 (Score 16)", "Heavy Lift Crane Vessel Weather Standby Rates", "1.4.1", "K. Solberg", "$1,800,000", "Flexible weather window with Heerema; monitor 7-day wave height."],
        ["R03 (Score 12)", "Subsea Tubular Steel Mill Delivery Lags", "1.2.1", "M. Berg", "$1,200,000", "Dual-source steel from European backup mills; hot-shot freight."],
        ["R04 (Score 9)", "Topside Lifting & Mating Mechanical Interference", "1.4.2", "T. Nygård", "$600,000", "Perform 3D laser scan trial fit prior to offshore mobilization."],
        ["R05 (Score 6)", "Structural Engineering AFC Interface Errors", "1.1.1", "H. Lindqvist", "$300,000", "Enforce 100% 3D CAD clash detection & DNV class verification."]
    ]
    for r_idx, r_data in enumerate(r_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = risk_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            if c_idx in [0, 4] and r_idx in [1, 2]:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # =========================================================================
    # SLIDE 10: VERTICAL SWIMLANE WBS BREAKDOWN
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide10, "9. Vertical Swimlane WBS Breakdown by Category", "Control Account Summary per Discipline ($26.50M BAC → $35.41M Outturn EAC)")

    wbs_table = slide10.shapes.add_table(6, 5, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)).table
    w_headers = ["WBS Category", "Baseline BAC", "Outturn EAC", "Status %", "Key Control Deliverables"]
    for i, h in enumerate(w_headers):
        cell = wbs_table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE

    w_data = [
        ["1.1 Engineering", "$3.20M", "$3.50M", "100% Complete", "1.1.1 Structural Detail ($1.8M), 1.1.2 Piping ($1.4M)"],
        ["1.2 Procurement", "$8.50M", "$9.10M", "92.5% Complete", "1.2.1 Tubular Steel ($5.0M), 1.2.2 Drilling Pumps ($3.5M)"],
        ["1.3 Yard Fabrication", "$9.80M", "$12.40M", "77.5% In Progress", "1.3.1 Verdal ($5.0M), 1.3.2 Egersund Mast Rework ($4.8M)"],
        ["1.4 Offshore Installation", "$5.00M", "$6.40M", "15% In Progress", "1.4.1 Heavy Lift ($3.2M), 1.4.2 Topside Mating ($1.8M)"],
        ["1.5 Commissioning", "$1.80M", "$2.00M", "0% Unstarted", "1.5.1 Hook-up & NDT ($1.0M), 1.5.2 Pre-Comm ($0.8M)"]
    ]
    for r_idx, r_data in enumerate(w_data, start=1):
        for c_idx, val in enumerate(r_data):
            cell = wbs_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if c_idx == 2 and r_idx in [3, 4]:
                p.font.bold = True
                p.font.color.rgb = C_RED

    # Save presentation
    output_pptx = os.path.join(prs_dir, "Drill_Tower_Executive_Steering_Presentation.pptx") if 'prs_dir' in locals() else os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "06_docs", "Drill_Tower_Executive_Steering_Presentation.pptx")
    prs.save(output_pptx)
    print(f"✅ Generated: 06_docs/Drill_Tower_Executive_Steering_Presentation.pptx")

if __name__ == "__main__":
    create_deck()
