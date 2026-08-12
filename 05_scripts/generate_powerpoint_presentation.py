"""
Generates a 16:9 Executive PowerPoint Presentation for the Offshore EPC Platform Drill Tower Project.
Uses python-pptx to build slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
DOCS_DIR = os.path.join(BASE_DIR, "06_docs")
PPTX_PATH = os.path.join(DOCS_DIR, "Drill_Tower_Project_Steering_Presentation.pptx")

# Color Palette (Edward Tufte Muted Corporate Palette)
DARK_BLUE = RGBColor(17, 24, 39)    # #111827
ACCENT_BLUE = RGBColor(37, 99, 235) # #2563EB
RAG_RED = RGBColor(220, 38, 38)     # #DC2626
RAG_AMBER = RGBColor(217, 119, 6)   # #D97706
RAG_GREEN = RGBColor(5, 150, 105)   # #059669
LIGHT_BG = RGBColor(249, 250, 251)  # #F9FAFB
TEXT_MUTED = RGBColor(107, 114, 128) # #6B7280

def add_header(slide, title_text, subtitle_text=""):
    """Adds a standardized top banner to a slide."""
    # Top bar background
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.color.rgb = DARK_BLUE
    
    tf = top_bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(209, 213, 219)

def build_presentation():
    print(f"--- Generating PowerPoint Presentation: {PPTX_PATH} ---")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BLUE
    bg1.line.color.rgb = DARK_BLUE
    
    tf1 = bg1.text_frame
    tf1.margin_left = Inches(1.0)
    tf1.margin_top = Inches(2.2)
    
    p1 = tf1.paragraphs[0]
    p1.text = "OFFSHORE EPC PLATFORM DRILL TOWER PROJECT"
    p1.font.bold = True
    p1.font.size = Pt(32)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Executive Steering Committee Status Briefing & Commercial Appraisal"
    p1_sub.font.size = Pt(18)
    p1_sub.font.color.rgb = ACCENT_BLUE
    
    p1_meta = tf1.add_paragraph()
    p1_meta.text = "\nStatus Date: August 31, 2026 (Month 8 / Status Week 36)\nAuthor: Frank Ellingsen, Lead Project Controller"
    p1_meta.font.size = Pt(13)
    p1_meta.font.color.rgb = RGBColor(156, 163, 175)

    # -------------------------------------------------------------------------
    # SLIDE 2: OVERALL HEALTH & EVM DASHBOARD
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Project Health & Core EVM Dashboard", "Month 8 (Aug 2026) Earned Value Performance Metrics")
    
    # Left Box: EVM Table
    table_shape = s2.shapes.add_table(7, 4, Inches(0.5), Inches(1.4), Inches(8.0), Inches(5.5))
    table = table_shape.table
    
    headers = ["Metric", "Description", "Value", "Status"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    evm_rows = [
        ["BAC", "Baseline Contract Budget", "$26,500,000", "Approved"],
        ["PV", "Planned Value Baseline", "$20,500,000", "77.4% Planned"],
        ["EV", "Earned Value Complete", "$17,540,000", "66.2% Complete"],
        ["AC", "Cumulative Actual Cost", "$23,440,000", "CRITICAL OVERRUN"],
        ["CPI", "Cost Efficiency Index", "0.7483", "Red ($0.75 / $1.00)"],
        ["SPI_t", "Earned Schedule Velocity", "0.9250", "Amber (-18.2 Days)"]
    ]
    for row_idx, r_data in enumerate(evm_rows, 1):
        for col_idx, val in enumerate(r_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if col_idx == 3 and "CRITICAL" in val:
                    p.font.color.rgb = RAG_RED
                    p.font.bold = True
                    
    # Right Box: Executive Callout Card
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.4), Inches(4.0), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(254, 242, 242)
    card.line.color.rgb = RAG_RED
    
    tf_card = card.text_frame
    tf_card.margin_left = Inches(0.3)
    tf_card.margin_top = Inches(0.3)
    
    pc = tf_card.paragraphs[0]
    pc.text = "🚨 EXECUTIVE ALERT"
    pc.font.bold = True
    pc.font.size = Pt(16)
    pc.font.color.rgb = RAG_RED
    
    p_body = tf_card.add_paragraph()
    p_body.text = "\n• Baseline Budget Depletion: Month 9 (Sep 2026)\n" \
                 "• Required Financing Line: +$8,913,604\n" \
                 "• Outturn Cost EAC: $35.41M (+33.6% Overrun)\n" \
                 "• Outturn Completion: Jan 31, 2027 (+31 Days Delay)\n" \
                 "• Primary Driver: Egersund Mast Rework ($2.40M) & Delay Overhead ($3.81M)"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = DARK_BLUE

    # -------------------------------------------------------------------------
    # SLIDE 3: COMMERCIAL FINANCIAL RATIOS
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Commercial Capital Budgeting & Appraisal Ratios", "Project Investment Return Metrics (10.0% WACC Hurdle Rate)")
    
    t_fin_shape = s3.shapes.add_table(9, 4, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    t_fin = t_fin_shape.table
    
    fin_headers = ["Commercial Metric", "Calculated Value", "Benchmark / Hurdle", "Commercial Evaluation"]
    for i, h in enumerate(fin_headers):
        cell = t_fin.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    fin_rows = [
        ["Net Present Value (NPV @ 10%)", "+$14,899,563", "$0.00 Hurdle", "Strong Positive Net Capital Value"],
        ["Internal Rate of Return (IRR)", "18.86%", "10.0% Hurdle Rate", "Outperforms Hurdle Target by +886 bps"],
        ["Simple Payback Period", "4.43 Years", "< 5.0 Years Target", "53.1 Months (May 2031) Payback"],
        ["Profitability Index (PI)", "1.42", "> 1.0 Ratio", "Generates $1.42 PV per $1.00 spent"],
        ["Total Simple ROI", "134.37%", "100.0% Base", "High 10-Year Cumulative Return"],
        ["Annualized ROI (CAGR)", "8.89%/Year", "5.0% Risk-Free", "Compounded Annual Growth Rate"],
        ["Gross Future Value (FV Yr 10)", "$130,499,397", "CAPEX Base", "Total Operating Inflow"],
        ["Net Future Value (NFV)", "+$38,645,628", "CAPEX Outturn", "Net Future Cash Surplus"]
    ]
    for row_idx, r_data in enumerate(fin_rows, 1):
        for col_idx, val in enumerate(r_data):
            cell = t_fin.cell(row_idx, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)

    # -------------------------------------------------------------------------
    # SLIDE 4: MONTE CARLO RISK SIMULATION (P90)
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Monte Carlo Risk Simulation (10,000 Iterations)", "Statistical Outturn Percentiles & Contingency Reserves")
    
    t_mc_shape = s4.shapes.add_table(6, 6, Inches(0.5), Inches(1.4), Inches(12.333), Inches(4.5))
    t_mc = t_mc_shape.table
    
    mc_headers = ["Percentile", "Confidence Level", "Outturn EAC", "Cost Overrun", "Completion Date", "Schedule Delay"]
    for i, h in enumerate(mc_headers):
        cell = t_mc.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(127, 29, 29)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    mc_rows = [
        ["P10", "10% Optimistic", "$32,444,302", "-$5,944,302", "Feb 13, 2027", "+43.1 Days"],
        ["P50", "50% Median", "$34,060,783", "-$7,560,783", "Feb 27, 2027", "+57.3 Days"],
        ["P80", "80% Standard Budget", "$35,195,026", "-$8,695,026", "Mar 09, 2027", "+67.3 Days"],
        ["P90", "90% High Confidence", "$35,815,202", "-$9,315,202", "Mar 14, 2027", "+72.5 Days"],
        ["P95", "95% Extreme Risk", "$36,272,986", "-$9,772,986", "Mar 18, 2027", "+76.6 Days"]
    ]
    for row_idx, r_data in enumerate(mc_rows, 1):
        for col_idx, val in enumerate(r_data):
            cell = t_mc.cell(row_idx, col_idx)
            cell.text = val
            if row_idx == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(254, 242, 242)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if row_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = RAG_RED
                    
    # Bottom Callout Box for P90
    card_mc = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.9))
    card_mc.fill.solid()
    card_mc.fill.fore_color.rgb = DARK_BLUE
    card_mc.line.color.rgb = DARK_BLUE
    tf_mc = card_mc.text_frame
    tf_mc.margin_left = Inches(0.3)
    tf_mc.margin_top = Inches(0.15)
    pmc = tf_mc.paragraphs[0]
    pmc.text = "P90 HIGH-CONFIDENCE TARGET: $35,815,202 Outturn Cost | March 14, 2027 Completion Date\n" \
               "Contingency Reserve Needed: +$401,598 above base EAC to achieve 90% financial confidence."
    pmc.font.bold = True
    pmc.font.size = Pt(11)
    pmc.font.color.rgb = RGBColor(255, 255, 255)

    # -------------------------------------------------------------------------
    # SLIDE 5: RISK MATRIX & ACTIVE RISK REGISTER
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "5x5 Executive Risk Matrix & Active Register", "Top Project Risks R01 - R05 with Financial Exposure")
    
    t_risk_shape = s5.shapes.add_table(6, 6, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    t_risk = t_risk_shape.table
    
    risk_headers = ["Risk ID", "Risk Title", "Severity", "Financial Exposure", "EMV Value", "Accountable CAM"]
    for i, h in enumerate(risk_headers):
        cell = t_risk.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    risk_rows = [
        ["R01", "Egersund Derrick Mast Yard Rework", "5x5 Critical", "$2,400,000", "$2,400,000", "Lars Hansen"],
        ["R02", "North Sea Weather Standby Delay", "4x4 High", "$1,800,000", "$720,000", "Erik Solberg"],
        ["R03", "Tubular Steel Price Inflation", "3x3 Medium", "$600,000", "$180,000", "Ingrid Berg"],
        ["R04", "DNV Structural Redesign Review", "2x3 Medium", "$300,000", "$60,000", "Geir Nilsen"],
        ["R05", "Offshore Hook-up Crew Bottleneck", "2x2 Low", "$200,000", "$40,000", "Bjørn Lie"]
    ]
    for row_idx, r_data in enumerate(risk_rows, 1):
        for col_idx, val in enumerate(r_data):
            cell = t_risk.cell(row_idx, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if col_idx == 2 and "Critical" in val:
                    p.font.color.rgb = RAG_RED
                    p.font.bold = True

    # -------------------------------------------------------------------------
    # SLIDE 6: COST VARIANCE WATERFALL BRIDGE
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Cost Variance Waterfall Bridge ($BAC to EAC)", "Step-by-Step Cost Overrun Progression across Control Accounts")
    
    t_wf_shape = s6.shapes.add_table(10, 5, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    t_wf = t_wf_shape.table
    
    wf_headers = ["Step", "Control Account / Variance Component", "Incremental Overrun", "Cumulative Cost", "% Contribution"]
    for i, h in enumerate(wf_headers):
        cell = t_wf.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    wf_rows = [
        ["0", "Original Baseline Budget (BAC)", "$26,500,000", "$26,500,000", "0.0%"],
        ["1", "Detail Structural Engineering Revisions", "+$300,000", "$26,800,000", "3.4%"],
        ["2", "Procurement Steel Price Inflation", "+$600,000", "$27,400,000", "6.7%"],
        ["3", "Verdal Sub-Structure Fabrication", "+$200,000", "$27,600,000", "2.2%"],
        ["4", "Egersund Derrick Mast Yard Rework", "+$2,400,000", "$30,000,000", "26.9%"],
        ["5", "Heavy Lift Vessel Weather Standby", "+$1,300,000", "$31,300,000", "14.6%"],
        ["6", "Offshore Topside Lifting & Mating", "+$100,000", "$31,400,000", "1.1%"],
        ["7", "Hook-up & Commissioning Crew", "+$200,000", "$31,600,000", "2.2%"],
        ["8", "Time Delay Overhead Spread (SPI_t)", "+$3,813,604", "$35,413,604", "42.8%"]
    ]
    for row_idx, r_data in enumerate(wf_rows, 1):
        for col_idx, val in enumerate(r_data):
            cell = t_wf.cell(row_idx, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                if row_idx in [5, 9]:
                    p.font.bold = True
                    p.font.color.rgb = RAG_RED

    # -------------------------------------------------------------------------
    # SLIDE 7: STEERING COMMITTEE DECISION MATRIX & ACTION PLAN
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Steering Committee Decision Matrix & Action Plan", "Required Executive Approvals & Project Recovery Directives")
    
    card_act = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    card_act.fill.solid()
    card_act.fill.fore_color.rgb = LIGHT_BG
    card_act.line.color.rgb = DARK_BLUE
    
    tf_act = card_act.text_frame
    tf_act.margin_left = Inches(0.5)
    tf_act.margin_top = Inches(0.4)
    
    p_act = tf_act.paragraphs[0]
    p_act.text = "RECOMMENDED STEERING COMMITTEE ACTIONS:"
    p_act.font.bold = True
    p_act.font.size = Pt(16)
    p_act.font.color.rgb = DARK_BLUE
    
    actions = [
        "\n1. AUTHORIZE +$8.91M OVERRUN FINANCING LINE:",
        "   • Executive Board approval required before September 15, 2026 to prevent Month 9 yard shutdown.",
        "\n2. ENFORCE FIXED-FEE CAP ON EGERSUND YARD BILLING:",
        "   • Transition remaining Derrick Mast assembly hours to a fixed-fee cap to halt further cost escalation.",
        "\n3. COMPRESS OFFSHORE HOOK-UP SCHEDULE:",
        "   • Recover 15 calendar days along the critical path to preserve the January 31, 2027 Commercial COD."
    ]
    for a in actions:
        pa = tf_act.add_paragraph()
        pa.text = a
        pa.font.size = Pt(12 if a.startswith("\n") else 11)
        if a.startswith("\n"):
            pa.font.bold = True
            pa.font.color.rgb = ACCENT_BLUE
        else:
            pa.font.color.rgb = DARK_BLUE
            
    prs.save(PPTX_PATH)
    print(f" [PASS] PowerPoint presentation generated successfully at: {PPTX_PATH}")

if __name__ == "__main__":
    build_presentation()
